import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: System Architecture
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Three-Tier Architecture of CityPulse"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Presentation Layer: Built with React.js. Provides dedicated interfaces for both standard Clients (Users) and the Admin Panel."

p = tf.add_paragraph()
p.text = "Application Layer: Powered by an Express.js API Server. It orchestrates several microservices:"
p.level = 0

p = tf.add_paragraph()
p.text = "Machine Learning Service: For automated image verification."
p.level = 1

p = tf.add_paragraph()
p.text = "Authentication Service: Secure user and admin login management."
p.level = 1

p = tf.add_paragraph()
p.text = "Notification Service: Real-time updates on report status."
p.level = 1

p = tf.add_paragraph()
p.text = "Data Layer: Robust data persistence utilizing a relational Database (Supabase/PostgreSQL) and Object Storage for user-uploaded images."
p.level = 0


# Slide 2: Data Models & Relationships
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "UML Class Diagram & Relationships"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Core Entities:"

p = tf.add_paragraph()
p.text = "User: Submits reports and earns ecoPoints."
p.level = 1

p = tf.add_paragraph()
p.text = "Admin: Verifies reports, manages users, and assigns severity."
p.level = 1

p = tf.add_paragraph()
p.text = "Central Entity - 'Report': Represents a civic issue. It maps to various essential sub-services:"
p.level = 0

p = tf.add_paragraph()
p.text = "Has 'Location' (Latitude, Longitude, Address)."
p.level = 1

p = tf.add_paragraph()
p.text = "Validated by 'MLService' (Automated AI image analysis)."
p.level = 1

p = tf.add_paragraph()
p.text = "Stored in 'DatabaseService'."
p.level = 1

p = tf.add_paragraph()
p.text = "Triggers 'NotificationManager' for real-time alerts."
p.level = 1


# Slide 3: Gamification & Rewards
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Gamification: What are Eco Coins?"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Concept: Eco Coins are the virtual currency of CityPulse, designed to incentivize active civic participation."

p = tf.add_paragraph()
p.text = "How to Earn: Users earn Eco Coins and XP by completing Daily Commissions (e.g., submitting garbage reports, checking the heatmap, viewing the leaderboard)."
p.level = 0

p = tf.add_paragraph()
p.text = "Progression: Eco Coins directly contribute to increasing your 'Adventurer Rank' and 'Login Streak,' transforming civic duties into a rewarding game."
p.level = 0

p = tf.add_paragraph()
p.text = "Future Use: Can be redeemed for exclusive profile badges, special titles, or real-world civic rewards."
p.level = 0


# Slide 4: Key Features
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Platform Features & Capabilities"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "AI-Powered Reporting: Upload a photo of garbage and the integrated AI automatically validates it and assigns a severity rating in real-time."

p = tf.add_paragraph()
p.text = "Interactive City Heatmap: A visual map showing active reports, enabling citizens and admins to identify high-priority zones at a glance."
p.level = 0

p = tf.add_paragraph()
p.text = "Gamified Dashboard & Leaderboard: A fully immersive 4-quadrant dashboard featuring daily quests, XP tracking, badges, and a global leaderboard to compete with other citizens."
p.level = 0

p = tf.add_paragraph()
p.text = "Comprehensive Admin Portal: A dedicated control panel for city officials to view reports, approve/reject submissions, promote users, and track analytics securely."
p.level = 0

# Save
prs.save('CityPulse_Presentation.pptx')
print("Presentation generated successfully at CityPulse_Presentation.pptx")
